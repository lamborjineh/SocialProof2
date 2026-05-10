"""
SocialProof — Module 0c: File Input
Extracts plain text from any uploaded document so the normal analysis pipeline
can process it regardless of format.

Supported formats:
  .pdf   — pdfplumber (layout-aware, handles columns/tables)
  .docx  — python-docx (paragraphs + tables)
  .pptx  — python-pptx (slide text frames)
  .html  — BeautifulSoup (strips tags, extracts body text)
  .txt   — plain UTF-8 read
  .json  — dumps values as readable text

Usage (called by routers/analyze.py):
    from pipeline.file_input import extract_text_from_file, is_file_support_available

    text = extract_text_from_file(file_bytes, filename="article.docx")

Notes:
  - Scanned/image-only PDFs return empty string — router should surface this
    so the user can re-submit as input_type='image' via OCR.
  - Password-protected PDFs fail gracefully and return empty string.
  - PDFs capped at MAX_PAGES to keep latency sane.
  - All handlers fail gracefully — missing optional library = clear log warning.
"""

import io
import json
import logging
from typing import Optional

logger = logging.getLogger("socialproof")

MAX_PAGES = 50  # PDF-only cap

# ── Format detection ──────────────────────────────────────────────────────────

EXTENSION_MAP = {
    ".pdf":  "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
    ".html": "html",
    ".htm":  "html",
    ".txt":  "txt",
    ".json": "json",
}

MIME_MAP = {
    "application/pdf":                                                  "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "text/html":                                                        "html",
    "text/plain":                                                       "txt",
    "application/json":                                                 "json",
}


def detect_format(filename: str = "", mime_type: str = "") -> Optional[str]:
    """Detect file format from filename extension or MIME type."""
    if filename:
        ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        if ext in EXTENSION_MAP:
            return EXTENSION_MAP[ext]
    if mime_type:
        base = mime_type.split(";")[0].strip().lower()
        if base in MIME_MAP:
            return MIME_MAP[base]
    return None


# ── Per-format extractors ─────────────────────────────────────────────────────

def _extract_pdf(file_bytes: bytes) -> Optional[str]:
    try:
        import pdfplumber

        pages_text: list[str] = []
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            total = len(pdf.pages)
            if total > MAX_PAGES:
                logger.warning(f"[File] PDF has {total} pages — capping at {MAX_PAGES}.")
            for page in pdf.pages[:MAX_PAGES]:
                try:
                    text = page.extract_text()
                    if text and text.strip():
                        pages_text.append(text.strip())
                except Exception as e:
                    logger.debug(f"[File] PDF page extraction failed: {e}")
                    continue

        combined = "\n\n".join(pages_text).strip()
        if combined:
            logger.info(f"[File] PDF: extracted {len(combined)} chars.")
            return combined

        logger.warning("[File] PDF: no text layer found — may be a scanned document.")
        return None

    except ImportError:
        logger.warning("[File] pdfplumber not installed. Run: pip install pdfplumber")
        return None
    except Exception as e:
        logger.warning(f"[File] PDF extraction failed: {e}")
        return None


def _extract_docx(file_bytes: bytes) -> Optional[str]:
    try:
        from docx import Document

        doc = Document(io.BytesIO(file_bytes))
        parts: list[str] = []

        for para in doc.paragraphs:
            t = para.text.strip()
            if t:
                parts.append(t)

        # Also pull text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                if row_text:
                    parts.append(row_text)

        combined = "\n\n".join(parts).strip()
        if combined:
            logger.info(f"[File] DOCX: extracted {len(combined)} chars.")
            return combined

        logger.warning("[File] DOCX: no text content found.")
        return None

    except ImportError:
        logger.warning("[File] python-docx not installed. Run: pip install python-docx")
        return None
    except Exception as e:
        logger.warning(f"[File] DOCX extraction failed: {e}")
        return None


def _extract_pptx(file_bytes: bytes) -> Optional[str]:
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(file_bytes))
        slides_text: list[str] = []

        for i, slide in enumerate(prs.slides, 1):
            slide_parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs).strip()
                        if t:
                            slide_parts.append(t)
            if slide_parts:
                slides_text.append(f"[Slide {i}]\n" + "\n".join(slide_parts))

        combined = "\n\n".join(slides_text).strip()
        if combined:
            logger.info(f"[File] PPTX: extracted {len(combined)} chars from {len(prs.slides)} slides.")
            return combined

        logger.warning("[File] PPTX: no text content found.")
        return None

    except ImportError:
        logger.warning("[File] python-pptx not installed. Run: pip install python-pptx")
        return None
    except Exception as e:
        logger.warning(f"[File] PPTX extraction failed: {e}")
        return None


def _extract_html(file_bytes: bytes) -> Optional[str]:
    try:
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(file_bytes.decode("utf-8", errors="ignore"), "html.parser")

        # Remove script, style, nav, footer noise
        for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
            tag.decompose()

        text = soup.get_text(separator="\n")
        lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
        combined = "\n".join(lines)

        if combined:
            logger.info(f"[File] HTML: extracted {len(combined)} chars.")
            return combined

        logger.warning("[File] HTML: no text content found.")
        return None

    except ImportError:
        logger.warning("[File] beautifulsoup4 not installed. Run: pip install beautifulsoup4")
        return None
    except Exception as e:
        logger.warning(f"[File] HTML extraction failed: {e}")
        return None


def _extract_txt(file_bytes: bytes) -> Optional[str]:
    try:
        text = file_bytes.decode("utf-8", errors="ignore").strip()
        if text:
            logger.info(f"[File] TXT: extracted {len(text)} chars.")
            return text
        logger.warning("[File] TXT: file was empty.")
        return None
    except Exception as e:
        logger.warning(f"[File] TXT extraction failed: {e}")
        return None


def _extract_json(file_bytes: bytes) -> Optional[str]:
    try:
        data = json.loads(file_bytes.decode("utf-8", errors="ignore"))

        def flatten(obj, depth=0) -> list[str]:
            parts = []
            if isinstance(obj, dict):
                for k, v in obj.items():
                    parts.append(f"{'  ' * depth}{k}:")
                    parts.extend(flatten(v, depth + 1))
            elif isinstance(obj, list):
                for item in obj:
                    parts.extend(flatten(item, depth))
            elif obj is not None:
                parts.append(f"{'  ' * depth}{obj}")
            return parts

        combined = "\n".join(flatten(data)).strip()
        if combined:
            logger.info(f"[File] JSON: extracted {len(combined)} chars.")
            return combined

        logger.warning("[File] JSON: no content found.")
        return None

    except Exception as e:
        logger.warning(f"[File] JSON extraction failed: {e}")
        return None


# ── Public interface ──────────────────────────────────────────────────────────

_EXTRACTORS = {
    "pdf":  _extract_pdf,
    "docx": _extract_docx,
    "pptx": _extract_pptx,
    "html": _extract_html,
    "txt":  _extract_txt,
    "json": _extract_json,
}

SUPPORTED_EXTENSIONS = list(EXTENSION_MAP.keys())  # [".pdf", ".docx", ...]
SUPPORTED_ACCEPT     = ",".join(SUPPORTED_EXTENSIONS)  # for HTML input accept=""


def extract_text_from_file(
    file_bytes: bytes,
    filename: str = "",
    mime_type: str = "",
) -> str:
    """
    Extract plain text from any supported file type.

    Args:
        file_bytes: Raw bytes of the uploaded file.
        filename:   Original filename (used for extension detection).
        mime_type:  MIME type string (fallback if filename has no extension).

    Returns:
        Extracted text string, or empty string on failure.
        Caller should check len(result) > 0 before proceeding.
    """
    if not file_bytes:
        return ""

    fmt = detect_format(filename, mime_type)
    if not fmt:
        logger.warning(
            f"[File] Could not detect format for filename='{filename}' "
            f"mime='{mime_type}'. Supported: {list(EXTENSION_MAP.keys())}"
        )
        return ""

    extractor = _EXTRACTORS.get(fmt)
    if not extractor:
        logger.warning(f"[File] No extractor registered for format '{fmt}'.")
        return ""

    result = extractor(file_bytes)
    return result or ""


def is_file_support_available() -> dict:
    """
    Health check — reports which extraction libraries are importable.
    Called by GET /health.
    """
    status = {}
    checks = {
        "pdfplumber":    "pdf",
        "docx":          "docx",
        "pptx":          "pptx",
        "bs4":           "html",
    }
    for lib, fmt in checks.items():
        try:
            __import__(lib)
            status[fmt] = True
        except ImportError:
            status[fmt] = False

    # txt and json need no external libs
    status["txt"]  = True
    status["json"] = True
    status["any"]  = any(status.values())
    return status
